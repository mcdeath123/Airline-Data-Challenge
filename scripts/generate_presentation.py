"""
generate_presentation.py

This script generates a stunning, premium PowerPoint presentation (.pptx)
for the Airline Data Challenge. It reflects the exact analytical findings,
financial models, and recommended portfolio, mirroring the modern, high-impact aesthetics
of the HTML slide deck. It inserts the generated visualization charts directly.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- STYLE DICTIONARY (DARK CYBER-PUNK / EXECUTIVE PREMIUM) ---
COLOR_BG = RGBColor(10, 14, 23)         # Sleek dark navy (#0A0E17)
COLOR_CARD = RGBColor(18, 26, 47)       # Premium glass card (#121A2F)
COLOR_BORDER = RGBColor(40, 52, 79)     # Soft subtle border (#28344F)
COLOR_PRIMARY = RGBColor(30, 61, 89)    # Classic primary blue (#1E3D59)
COLOR_ACCENT = RGBColor(255, 110, 64)   # Sunrise orange accent (#FF6E40)
COLOR_GOLD = RGBColor(255, 193, 59)     # Golden highlight (#FFC13B)
COLOR_TEXT_LIGHT = RGBColor(243, 244, 246) # White text (#F3F4F6)
COLOR_TEXT_MUTED = RGBColor(156, 163, 175) # Muted grey (#9CA3AF)

FONT_TITLE = 'Trebuchet MS'
FONT_BODY = 'Segoe UI'

def create_deck():
    print("Initializing PowerPoint Slide Deck...")
    prs = Presentation()
    
    # Set standard widescreen 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ------------------ HELPER FUNCTIONS ------------------
    
    def apply_dark_bg(slide):
        """Creates a solid background covering the entire slide."""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, subtitle_num, title_text):
        """Adds a standardized premium header zone to a slide."""
        # Header Text Box
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        
        # Subtitle / Section label
        p_sub = tf.paragraphs[0]
        p_sub.text = subtitle_num.upper()
        p_sub.font.name = FONT_TITLE
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_ACCENT
        p_sub.space_after = Pt(2)
        
        # Main Title
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_LIGHT
        
        # Decorative divider line under header
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BORDER
        line.line.fill.background()

    def add_card(slide, left, top, width, height, border_color=COLOR_BORDER):
        """Creates a visual glassmorphism card component with an optional accent border."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    def add_bullet_list(slide, left, top, width, height, bullets, font_size=13):
        """Adds a clean, well-spaced list of bullet points with key terms bolded."""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(12)
            p.line_spacing = 1.15
            
            # Custom bullet symbol "✦ "
            run_bullet = p.add_run()
            run_bullet.text = "✦  "
            run_bullet.font.name = FONT_BODY
            run_bullet.font.size = Pt(font_size)
            run_bullet.font.bold = True
            run_bullet.font.color.rgb = COLOR_ACCENT
            
            # Parse double asterisks for bolding key concepts
            parts = bullet.split("**")
            for i, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.name = FONT_BODY
                run.font.size = Pt(font_size)
                run.font.color.rgb = COLOR_TEXT_LIGHT if i % 2 == 1 else COLOR_TEXT_MUTED
                if i % 2 == 1:
                    run.font.bold = True

    def add_metric_card(slide, left, top, width, height, value_text, label_text):
        """Adds a visual metric card with a large golden figure."""
        add_card(slide, left, top, width, height, border_color=COLOR_PRIMARY)
        
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.15)
        
        # Value paragraph
        p_val = tf.paragraphs[0]
        p_val.alignment = PP_ALIGN.CENTER
        p_val.text = value_text
        p_val.font.name = FONT_TITLE
        p_val.font.size = Pt(28)
        p_val.font.bold = True
        p_val.font.color.rgb = COLOR_GOLD
        p_val.space_after = Pt(4)
        
        # Label paragraph
        p_lbl = tf.add_paragraph()
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.text = label_text.upper()
        p_lbl.font.name = FONT_BODY
        p_lbl.font.size = Pt(9)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_TEXT_MUTED

    def add_table_to_slide(slide, left, top, width, height, headers, rows, highlight_first_row=False, col_widths=None):
        """Draws a beautiful, premium table with styled cells."""
        cols_count = len(headers)
        rows_count = len(rows) + 1 # Include header row
        
        table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
        table = table_shape.table
        
        # Apply custom column widths if provided
        if col_widths and len(col_widths) == cols_count:
            for col_idx, col_w in enumerate(col_widths):
                table.columns[col_idx].width = Inches(col_w)
        
        # Style Headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            
            p = cell.text_frame.paragraphs[0]
            p.text = header
            p.font.name = FONT_TITLE
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = COLOR_ACCENT
            p.alignment = PP_ALIGN.LEFT
            
        # Style Data Rows
        for row_idx, row_data in enumerate(rows):
            is_highlight = (row_idx == 0 and highlight_first_row)
            for col_idx, val in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.fill.solid()
                if is_highlight:
                    cell.fill.fore_color.rgb = RGBColor(30, 41, 59) # Deeper blue-slate for highlight
                else:
                    cell.fill.fore_color.rgb = COLOR_BG
                    
                p = cell.text_frame.paragraphs[0]
                p.text = str(val)
                p.font.name = FONT_BODY
                p.font.size = Pt(10)
                p.font.color.rgb = COLOR_TEXT_LIGHT if (is_highlight or col_idx == 0) else COLOR_TEXT_MUTED
                if is_highlight or col_idx == 0:
                    p.font.bold = True
                p.alignment = PP_ALIGN.LEFT

    def add_visual_chart(slide, left, top, width, height, image_path):
        """Inserts a chart and wraps it in a sleek container block, preserving aspect ratio."""
        # Draw background backing card for the image
        add_card(slide, left, top, width, height, border_color=COLOR_BORDER)
        
        # Check if file exists, then insert it
        if os.path.exists(image_path):
            padding = Inches(0.12)
            max_w = width - (padding * 2)
            max_h = height - (padding * 2)
            
            # Load picture at its default size first
            pic = slide.shapes.add_picture(image_path, left, top)
            
            # Calculate natural aspect ratio
            aspect_ratio = pic.width / pic.height
            
            # Re-scale to fit within max_w and max_h boundaries
            if max_w / aspect_ratio <= max_h:
                pic.width = max_w
                pic.height = int(max_w / aspect_ratio)
            else:
                pic.height = max_h
                pic.width = int(max_h * aspect_ratio)
                
            # Center the picture within the card container
            pic.left = left + int((width - pic.width) / 2)
            pic.top = top + int((height - pic.height) / 2)
        else:
            # Placeholder text if image is missing
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"[Chart Image Missing:\n{image_path}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT_BODY
            p.font.color.rgb = COLOR_ACCENT
            p.font.size = Pt(12)

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (STUNNING WELCOME)
    # =========================================================================
    print("Building Slide 1: Title...")
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide1)
    
    # Abstract diagonal geometric element (Accent Line)
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), prs.slide_height)
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()
    
    # Title Block
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    # Badge
    p_badge = tf1.paragraphs[0]
    p_badge.text = "BUSINESS INTELLIGENCE PITCH"
    p_badge.font.name = FONT_TITLE
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_ACCENT
    p_badge.space_after = Pt(16)
    
    # Large Title
    p_title = tf1.add_paragraph()
    p_title.text = "US MARKET ENTRY STRATEGY"
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_LIGHT
    p_title.space_after = Pt(8)
    
    # Subtitle
    p_sub = tf1.add_paragraph()
    p_sub.text = "1Q2019 Aviation Market Analysis & Route Investment Recommendations"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sub.space_after = Pt(50)
    
    # Metadata footer
    p_meta = tf1.add_paragraph()
    p_meta.text = "Airline Data Challenge  |  Corporate Strategy  |  May 2026"
    p_meta.font.name = FONT_BODY
    p_meta.font.size = Pt(11)
    p_meta.font.color.rgb = COLOR_ACCENT
    
    # =========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY
    # =========================================================================
    print("Building Slide 2: Executive Summary...")
    slide2 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide2)
    add_header(slide2, "01 / Strategic Blueprint", "Executive Summary")
    
    # Left Column: Bullet points
    summary_bullets = [
        "**The Challenge**: Establish a highly competitive domestic airline in the United States, utilizing an initial fleet of **5 aircraft**.",
        "**Upfront Aircraft Cost**: A substantial upfront capital expenditure of **$90 million per aircraft** ($450 million portfolio total).",
        "**Core Motto**: *'On time, for you.'* Our brand is anchored on punctuality, which demands incorporating real-world delays into our financial risk evaluations.",
        "**Recommended Solution**: Operate a diversified network of 5 highly liquid transcontinental and regional routes, generating **$502.2M in profit**."
    ]
    add_bullet_list(slide2, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0), summary_bullets)
    
    # Right Column: Big Metrics Grid
    # Row 1 metrics
    add_metric_card(slide2, Inches(7.4), Inches(1.8), Inches(2.4), Inches(1.5), "$450M", "Fleet Capital CapEx")
    add_metric_card(slide2, Inches(10.133), Inches(1.8), Inches(2.4), Inches(1.5), "5", "Aircraft Active")
    
    # Row 2 metrics
    add_metric_card(slide2, Inches(7.4), Inches(3.7), Inches(2.4), Inches(1.5), "$502.2M", "Portfolio Net Profit")
    add_metric_card(slide2, Inches(10.133), Inches(3.7), Inches(2.4), Inches(1.5), "65.5%", "Average Occupancy")
    
    # Strategic Note Card below metrics
    add_card(slide2, Inches(7.4), Inches(5.6), Inches(5.133), Inches(1.2), border_color=COLOR_ACCENT)
    note_box = slide2.shapes.add_textbox(Inches(7.5), Inches(5.65), Inches(4.933), Inches(1.1))
    ntf = note_box.text_frame
    ntf.word_wrap = True
    np = ntf.paragraphs[0]
    np.text = "STRATEGIC MANDATE:"
    np.font.name = FONT_TITLE
    np.font.size = Pt(10)
    np.font.bold = True
    np.font.color.rgb = COLOR_ACCENT
    np.space_after = Pt(2)
    np2 = ntf.add_paragraph()
    np2.text = "Target high-occupancy transcontinental channels and extreme-yield short hauls to lock in market dominance and accelerate capital amortization."
    np2.font.name = FONT_BODY
    np2.font.size = Pt(10)
    np2.font.color.rgb = COLOR_TEXT_LIGHT
    
    # =========================================================================
    # SLIDE 3: DATA QUALITY & CLEANING
    # =========================================================================
    print("Building Slide 3: Data Quality...")
    slide3 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide3)
    add_header(slide3, "02 / Data Management", "Data Quality & Imputation Rules")
    
    # Grid of 4 card components for the 4 quality insights
    # Row 1 Card 1
    add_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.4), border_color=COLOR_ACCENT)
    c1_box = slide3.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(5.3), Inches(2.1))
    c1_tf = c1_box.text_frame
    c1_tf.word_wrap = True
    p = c1_tf.paragraphs[0]
    p.text = "1. Reward Ticket & Security Fee Clean"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.space_after = Pt(6)
    p2 = c1_tf.add_paragraph()
    p2.text = "Over 5% of fares were under $11 (the standard Sept 11 security fee). Passenger yield is distorted if included. Filtered out all tickets under $50 (reward tickets) and over $5,000 (corrupted segment rates) to isolate commercial fares."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.line_spacing = 1.1
    
    # Row 1 Card 2
    add_card(slide3, Inches(6.933), Inches(1.8), Inches(5.6), Inches(2.4))
    c2_box = slide3.shapes.add_textbox(Inches(7.083), Inches(1.95), Inches(5.3), Inches(2.1))
    c2_tf = c2_box.text_frame
    c2_tf.word_wrap = True
    p = c2_tf.paragraphs[0]
    p.text = "2. Missing Distance Imputation"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(6)
    p2 = c2_tf.add_paragraph()
    p2.text = "We encountered 2,700 records containing null distances or **** characters. Imputed using the median distance of all non-null flights on the exact same directed route (e.g. LAX->JFK). Protects distance-based cost calculations ($9.18/mi)."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.line_spacing = 1.1
 
    # Row 2 Card 1
    add_card(slide3, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.4))
    c3_box = slide3.shapes.add_textbox(Inches(0.95), Inches(4.65), Inches(5.3), Inches(2.1))
    c3_tf = c3_box.text_frame
    c3_tf.word_wrap = True
    p = c3_tf.paragraphs[0]
    p.text = "3. Diverted Flight Delay Management"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.space_after = Pt(6)
    p2 = c3_tf.add_paragraph()
    p2.text = "Identified 4,377 non-cancelled flights that had valid departure delays but null arrival delays (representing diverted flights). Imputed arrival delay using the departure delay value to preserve volume and delay metrics."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.line_spacing = 1.1
 
    # Row 2 Card 2
    add_card(slide3, Inches(6.933), Inches(4.5), Inches(5.6), Inches(2.4))
    c4_box = slide3.shapes.add_textbox(Inches(7.083), Inches(4.65), Inches(5.3), Inches(2.1))
    c4_tf = c4_box.text_frame
    c4_tf.word_wrap = True
    p = c4_tf.paragraphs[0]
    p.text = "4. Mixed-Type Format Clean"
    p.font.name = FONT_TITLE
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.space_after = Pt(6)
    p2 = c4_tf.add_paragraph()
    p2.text = "Raw datasets contained massive formatting anomalies like '$$$' in air times and '$ 200.00' string values in ticket fares. Designed a high-speed cleaning pipeline that stripped currency symbols and coerced data to clean numeric values."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.line_spacing = 1.1

    # =========================================================================
    # SLIDE 4: DATA ARCHITECTURE & ENGINEERED FIELD METADATA [NEW SLIDE]
    # =========================================================================
    print("Building Slide 4: Data Dictionary & Metadata...")
    slide4 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide4)
    add_header(slide4, "03 / Data Engineering", "Engineered Field Formulas & Metadata Dictionary")
    
    # 6-Card Grid layout for Metadata formulas
    meta_cards = [
        {
            "title": "PASSENGERS",
            "formula": "200 * Occupancy Rate",
            "desc": "Expected passenger load per leg. Fixed mainline capacity scaled by real Flights occupancy data.",
            "color": COLOR_ACCENT
        },
        {
            "title": "REVENUE_TICKET",
            "formula": "Passengers * (Avg Fare / 2)",
            "desc": "Isolates flight leg fare yield. divided by 2 because joined Average Fare represents a round trip.",
            "color": COLOR_GOLD
        },
        {
            "title": "REVENUE_BAGGAGE",
            "formula": "Passengers * 50% * $35.00",
            "desc": "Ancillary baggage revenue per leg. Assumes a $35 checked bag fee at a conservative 50% check rate.",
            "color": COLOR_GOLD
        },
        {
            "title": "COST_DISTANCE",
            "formula": "Distance * $9.18 / mile",
            "desc": "Variable operating costs: fuel, crew, maintenance ($8.00/mi) + depreciation, insurance ($1.18/mi).",
            "color": COLOR_PRIMARY
        },
        {
            "title": "COST_AIRPORT",
            "formula": "Large: $10,000 | Medium: $5,000",
            "desc": "Fixed destination airport landing operational fee. Evaluated per leg using Airport_Codes.csv sizes.",
            "color": COLOR_PRIMARY
        },
        {
            "title": "COST_DELAY",
            "formula": "(Delay - 15) * $75 / minute",
            "desc": "Punctuality penalty. Tracks departure/arrival delays over 15 minutes at $75/min in added cost.",
            "color": COLOR_ACCENT
        }
    ]
    
    m_width = Inches(3.64)
    m_height = Inches(2.3)
    m_gap_x = Inches(0.4)
    m_gap_y = Inches(0.3)
    m_start_left = Inches(0.8)
    m_start_top = Inches(1.8)
    
    for idx, m_data in enumerate(meta_cards):
        col = idx % 3
        row = idx // 3
        
        m_left = m_start_left + (col * (m_width + m_gap_x))
        m_top = m_start_top + (row * (m_height + m_gap_y))
        
        # Add Card
        add_card(slide4, m_left, m_top, m_width, m_height, border_color=m_data["color"])
        
        # Text Block
        box = slide4.shapes.add_textbox(m_left + Inches(0.15), m_top + Inches(0.15), m_width - Inches(0.3), m_height - Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = True
        
        # Title
        p_title = tf.paragraphs[0]
        p_title.text = m_data["title"]
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = m_data["color"]
        p_title.space_after = Pt(2)
        
        # Formula
        p_form = tf.add_paragraph()
        p_form.text = f"Formula: {m_data['formula']}"
        p_form.font.name = FONT_BODY
        p_form.font.size = Pt(10)
        p_form.font.bold = True
        p_form.font.color.rgb = COLOR_TEXT_LIGHT
        p_form.space_after = Pt(6)
        
        # Description
        p_desc = tf.add_paragraph()
        p_desc.text = m_data["desc"]
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.line_spacing = 1.1
    
    # =========================================================================
    # SLIDE 5: BUSIEST ROUND-TRIP ROUTES
    # =========================================================================
    print("Building Slide 5: Busiest Routes...")
    slide5 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide5)
    add_header(slide5, "04 / Market Demand", "Busiest Round-Trip Routes")
    
    # Left Column: Bullet points and Table
    bullets_busiest = [
        "**Volume Indicated Liquidity**: High flight volume indicates deep, sustained commercial demand between cities.",
        "**Market Constraints**: However, high-volume routes are extremely competitive, leading to severe yield compression.",
        "**Transcontinental Standout**: **JFK-LAX** is the only transcontinental route in the top 5 busiest routes, showcasing its massive scale."
    ]
    add_bullet_list(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.0), bullets_busiest, font_size=11)
    
    # Top 5 Busiest Table
    headers_busiest = ["Rank", "Route", "Round Trips", "Avg Distance"]
    rows_busiest = [
        ["1", "LAX - SFO", "4,170.0", "337 mi"],
        ["2", "LGA - ORD", "3,578.0", "733 mi"],
        ["3", "LAS - LAX", "3,255.5", "236 mi"],
        ["4", "JFK - LAX (Rec)", "3,160.0", "2,475 mi"],
        ["5", "LAX - SEA", "2,499.5", "954 mi"]
    ]
    add_table_to_slide(slide5, Inches(0.8), Inches(4.0), Inches(5.6), Inches(2.8), headers_busiest, rows_busiest, highlight_first_row=False, col_widths=[0.7, 1.4, 1.7, 1.8])
    
    # Right Column: Chart 1
    add_visual_chart(slide5, Inches(6.8), Inches(1.8), Inches(5.733), Inches(5.0), "visualizations/chart_1_busiest_routes.png")
    
    # =========================================================================
    # SLIDE 6: MOST PROFITABLE ROUND-TRIP ROUTES
    # =========================================================================
    print("Building Slide 6: Profitable Routes...")
    slide6 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide6)
    add_header(slide6, "05 / Financial Assessment", "Most Profitable Round-Trip Routes")
    
    # Left Column: Bullet points and Table
    bullets_profitable = [
        "**Scale Matters**: Profits are calculated at the individual flight leg level by mapping fare average yield against physical costs.",
        "**Transcon Dominance**: Long-haul transcontinental flights sweep the top 3 spots, fueled by high ticket fares.",
        "**Regional Outlier**: **CLT-FLO** appears at #8 due to regional monopoly pricing, but it represents a tiny regional market."
    ]
    add_bullet_list(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.0), bullets_profitable, font_size=11)
    
    # Top 5 Profitable Table
    headers_profitable = ["Rank", "Route", "Q1 Net Profit", "Avg Fare", "Round Trips"]
    rows_profitable = [
        ["1", "JFK - LAX (Rec)", "$195.6M", "$964.35", "3,160.0"],
        ["2", "EWR - SFO (Rec)", "$82.8M", "$1,050.82", "1,212.0"],
        ["3", "JFK - SFO (Rec)", "$82.8M", "$860.38", "1,860.5"],
        ["4", "DCA - ORD (Rec)", "$72.4M", "$535.33", "1,847.5"],
        ["5", "ATL - CLT (Rec)", "$68.6M", "$508.35", "1,538.0"]
    ]
    add_table_to_slide(slide6, Inches(0.8), Inches(4.0), Inches(5.6), Inches(2.8), headers_profitable, rows_profitable, highlight_first_row=False, col_widths=[0.6, 1.4, 1.3, 1.0, 1.3])
    
    # Right Column: Chart 2
    add_visual_chart(slide6, Inches(6.8), Inches(1.8), Inches(5.733), Inches(5.0), "visualizations/chart_2_profitable_routes.png")

    # =========================================================================
    # SLIDE 7: RECOMMENDED PORTFOLIO (5 POWERFUL SLIDES IN 1 GRID)
    # =========================================================================
    print("Building Slide 7: Recommended Portfolio...")
    slide7 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide7)
    add_header(slide7, "06 / Fleet Investment", "Recommended Route Portfolio & Full Airport Identification")
    
    rec_cards_data = [
        {
            "route": "JFK-LAX",
            "desc": "New York John F. Kennedy Int'l (JFK) to Los Angeles Int'l (LAX). Undisputed premium transcon corridor.",
            "color": COLOR_ACCENT,
            "stats": [("Net Profit", "$195.6M"), ("Avg Fare", "$964.35"), ("RT Flights", "3,160.0")]
        },
        {
            "route": "EWR-SFO",
            "desc": "Newark Liberty Int'l (EWR) to San Francisco Int'l (SFO). Tech-finance corridor. Highest U.S. average fare.",
            "color": COLOR_GOLD,
            "stats": [("Net Profit", "$82.8M"), ("Avg Fare", "$1,050.82"), ("RT Flights", "1,212.0")]
        },
        {
            "route": "JFK-SFO",
            "desc": "New York JFK to San Francisco Int'l (SFO). Pairs with JFK-LAX & EWR-SFO for maximum network capture.",
            "color": RGBColor(0, 168, 204),
            "stats": [("Net Profit", "$82.8M"), ("Avg Fare", "$860.38"), ("RT Flights", "1,860.5")]
        },
        {
            "route": "DCA-ORD",
            "desc": "Washington Reagan National (DCA) to Chicago O'Hare Int'l (ORD). High-volume corporate business trunk.",
            "color": RGBColor(12, 123, 147),
            "stats": [("Net Profit", "$72.4M"), ("Avg Fare", "$535.33"), ("RT Flights", "1,847.5")]
        },
        {
            "route": "ATL-CLT",
            "desc": "Hartsfield-Jackson Atlanta (ATL) to Charlotte Douglas (CLT). Low-cost regional southeastern powerhouse.",
            "color": RGBColor(23, 88, 115),
            "stats": [("Net Profit", "$68.6M"), ("Avg Fare", "$508.35"), ("RT Flights", "1,538.0")]
        }
    ]
    
    card_width = Inches(2.15)
    card_gap = Inches(0.24)
    start_left = Inches(0.8)
    card_top = Inches(1.8)
    card_height = Inches(4.5)
    
    for idx, c_data in enumerate(rec_cards_data):
        c_left = start_left + (idx * (card_width + card_gap))
        
        # Add Card Base Shape
        add_card(slide7, c_left, card_top, card_width, card_height, border_color=c_data["color"])
        
        # Color accent header bar on top of card
        accent_bar = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, card_top, card_width, Inches(0.12))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = c_data["color"]
        accent_bar.line.fill.background()
        
        # Content Box
        c_box = slide7.shapes.add_textbox(c_left + Inches(0.1), card_top + Inches(0.2), card_width - Inches(0.2), card_height - Inches(0.3))
        c_tf = c_box.text_frame
        c_tf.word_wrap = True
        
        # Route title
        p_route = c_tf.paragraphs[0]
        p_route.text = c_data["route"]
        p_route.font.name = FONT_TITLE
        p_route.font.size = Pt(20)
        p_route.font.bold = True
        p_route.font.color.rgb = COLOR_TEXT_LIGHT
        p_route.space_after = Pt(10)
        
        # Description
        p_desc = c_tf.add_paragraph()
        p_desc.text = c_data["desc"]
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.space_after = Pt(14)
        p_desc.line_spacing = 1.1
        
        # Key metrics list
        for label, val in c_data["stats"]:
            p_stat = c_tf.add_paragraph()
            p_stat.space_after = Pt(3)
            
            # Label
            lbl_run = p_stat.add_run()
            lbl_run.text = f"{label}:\n"
            lbl_run.font.name = FONT_BODY
            lbl_run.font.size = Pt(9)
            lbl_run.font.color.rgb = COLOR_TEXT_MUTED
            
            # Value
            val_run = p_stat.add_run()
            val_run.text = val
            val_run.font.name = FONT_TITLE
            val_run.font.size = Pt(13)
            val_run.font.bold = True
            val_run.font.color.rgb = COLOR_GOLD
 
    # Strategic caution callout below cards
    add_card(slide7, Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.7), border_color=COLOR_ACCENT)
    caution_box = slide7.shapes.add_textbox(Inches(0.95), Inches(6.52), Inches(11.433), Inches(0.6))
    ctf = caution_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.space_after = Pt(0)
    
    c_lbl = cp.add_run()
    c_lbl.text = "CRITICAL STRATEGIC FOCUS: "
    c_lbl.font.name = FONT_TITLE
    c_lbl.font.bold = True
    c_lbl.font.size = Pt(10)
    c_lbl.font.color.rgb = COLOR_ACCENT
    
    c_text = cp.add_run()
    c_text.text = "Avoided regional monopoly CLT-FLO despite extreme raw fares ($1,999). Operating 200-seat aircraft in a small 3-flight-per-day market is physically unviable and would collapse prices. Mainline assets require high-density, deep markets."
    c_text.font.name = FONT_BODY
    c_text.font.size = Pt(9.5)
    c_text.font.color.rgb = COLOR_TEXT_LIGHT
    
    # =========================================================================
    # SLIDE 8: CAPITAL BREAKEVEN ANALYSIS
    # =========================================================================
    print("Building Slide 8: Breakeven Analysis...")
    slide8 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide8)
    add_header(slide8, "07 / Financial Breakeven", "Capital Breakeven Analysis ($90M Aircraft)")
    
    # Left Column: Bullet points and Table
    bullets_breakeven = [
        "**Upfront Investment**: Each dedicated aircraft requires a **$90,000,000** capital expenditure.",
        "**Breakeven Speed**: **EWR-SFO** offers the fastest flight amortization (1,317 RTs) due to its massive $68k profit per round trip.",
        "**Days to Breakeven**: Based on capturing average Q1 daily volumes, aircraft cost is fully amortized in under 120 days."
    ]
    add_bullet_list(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.0), bullets_breakeven, font_size=11)
    
    # Breakeven Table
    headers_breakeven = ["Route", "Profit / RT", "Breakeven RTs", "Projected Days"]
    rows_breakeven = [
        ["EWR - SFO", "$68,324", "1,317 RTs", "98 Days"],
        ["JFK - LAX", "$61,902", "1,454 RTs", "41 Days"],
        ["ATL - CLT", "$44,602", "2,018 RTs", "118 Days"],
        ["JFK - SFO", "$44,503", "2,022 RTs", "98 Days"],
        ["DCA - ORD", "$39,213", "2,295 RTs", "112 Days"]
    ]
    add_table_to_slide(slide8, Inches(0.8), Inches(4.0), Inches(5.6), Inches(2.8), headers_breakeven, rows_breakeven, highlight_first_row=True, col_widths=[1.1, 1.2, 1.7, 1.6])
    
    # Right Column: Chart 3
    add_visual_chart(slide8, Inches(6.8), Inches(1.8), Inches(5.733), Inches(5.0), "visualizations/chart_3_breakeven_analysis.png")
 
    # =========================================================================
    # SLIDE 9: PROFITABILITY DRIVERS
    # =========================================================================
    print("Building Slide 9: Profitability Drivers...")
    slide9 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide9)
    add_header(slide9, "08 / Business Drivers", "Recommended Routes: Profitability Drivers")
    
    # Left Column: Bullet points explaining variables
    bullets_drivers = [
        "**Dual Engines of Profit**: High load factor (occupancy rate) and premium fares are the core engines of profitability.",
        "**High Fare, High Yield**: **EWR-SFO** and **JFK-LAX** occupy the sweet spot of massive premium transcontinental yields.",
        "**Short Haul Capacity**: **ATL-CLT** bypasses normal distance limits with high fares ($508) and tiny operating costs ($2,074/leg).",
        "**Network Leverage**: Operating multiple endpoints from New York (JFK/EWR) to San Francisco/Los Angeles maximizes aircraft routing flexibility and capture."
    ]
    add_bullet_list(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), bullets_drivers)
    
    # Right Column: Chart 4
    add_visual_chart(slide9, Inches(6.8), Inches(1.8), Inches(5.733), Inches(5.0), "visualizations/chart_4_profitability_drivers.png")

    # =========================================================================
    # SLIDE 10: OPERATIONAL DELAY & BRAND PROMISE ALIGNMENT [NEW SLIDE]
    # =========================================================================
    print("Building Slide 10: Operational Delay Analysis...")
    slide10 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide10)
    add_header(slide10, "09 / Operational Risk", "On-Time Performance & Operational Delay Risk")
    
    # Left Column: Delay Bullet Points
    delay_bullets = [
        "**Brand Promise**: Our brand motto is *'On time, for you'*. Minimizing delays is not just a marketing slogan—it directly protects our high-yield margins.",
        "**Delay Cost Penalties**: Delays over 15 minutes incur heavy crew and gate penalties of **$75.00 per minute** for departure and arrival.",
        "**Newark (EWR) Turnaround Focus**: **EWR-SFO** has the highest average delay penalty (**$5,720 per round trip**). We must implement tight ground-turn execution and dedicated gates at Newark to mitigate this risk.",
        "**ATL-CLT Stability**: The regional short haul **ATL-CLT** is highly stable, with a tiny average delay penalty of only **$1,762 per round trip**."
    ]
    add_bullet_list(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), delay_bullets, font_size=11)
    
    # Right Column: Custom Delay Risk Table
    headers_delays = ["Route", "Q1 Delays", "Avg Delay / RT", "Risk Assessment"]
    rows_delays = [
        ["EWR - SFO", "$6.93M", "$5,720 / RT", "HIGH RISK (Turnaround Focus)"],
        ["JFK - SFO", "$8.17M", "$4,391 / RT", "MODERATE RISK (JFK Congestion)"],
        ["DCA - ORD", "$6.39M", "$3,456 / RT", "MODERATE RISK (Weather Slots)"],
        ["JFK - LAX", "$8.18M", "$2,588 / RT", "LOW-MODERATE RISK (Dense Trunk)"],
        ["ATL - CLT", "$2.71M", "$1,762 / RT", "LOW RISK (High Resilience)"]
    ]
    add_table_to_slide(slide10, Inches(6.8), Inches(1.8), Inches(5.733), Inches(3.0), headers_delays, rows_delays, highlight_first_row=True, col_widths=[1.1, 1.1, 1.4, 2.133])
    
    # Risk Mitigation Card below table
    add_card(slide10, Inches(6.8), Inches(5.2), Inches(5.733), Inches(1.5), border_color=COLOR_ACCENT)
    mit_box = slide10.shapes.add_textbox(Inches(6.95), Inches(5.25), Inches(5.433), Inches(1.4))
    mit_tf = mit_box.text_frame
    mit_tf.word_wrap = True
    mit_p = mit_tf.paragraphs[0]
    mit_p.text = "OPERATIONAL MITIGATION PLAN:"
    mit_p.font.name = FONT_TITLE
    mit_p.font.size = Pt(10)
    mit_p.font.bold = True
    mit_p.font.color.rgb = COLOR_ACCENT
    mit_p.space_after = Pt(2)
    mit_p2 = mit_tf.add_paragraph()
    mit_p2.text = "Establish standby crews and secondary routing options for our premium transcontinental routes (EWR/JFK to SFO) to absorb localized hub congestion and protect our customer on-time experience."
    mit_p2.font.name = FONT_BODY
    mit_p2.font.size = Pt(9.5)
    mit_p2.font.color.rgb = COLOR_TEXT_LIGHT
    mit_p2.line_spacing = 1.1

    # =========================================================================
    # SLIDE 11: STRATEGIC KEY PERFORMANCE INDICATORS (KPIs)
    # =========================================================================
    print("Building Slide 11: Strategic KPIs...")
    slide11 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide11)
    add_header(slide11, "10 / Operations Control", "Strategic Key Performance Indicators (KPIs)")
    
    kpi_cards = [
        {
            "num": "01",
            "title": "On-Time Performance (OTP)",
            "desc": "Track D15/A15 departure and arrival delay percentages.\n\n**Bottom Line**: Since delays over 15 minutes cost **$75 per minute**, maintaining a 90%+ OTP directly protects our operating margins."
        },
        {
            "num": "02",
            "title": "Route Load Factor",
            "desc": "Percentage of the 200 seats filled by paying passengers.\n\n**Bottom Line**: Occupancy is the core scaler of revenue. An extra 1% load factor on JFK-LAX yields **$1.5M/quarter** in high-margin revenues."
        },
        {
            "num": "03",
            "title": "Yield/Seat Mile (YASM)",
            "desc": "Total leg ticket and bag revenue divided by Available Seat Miles (Seats * Distance).\n\n**Bottom Line**: Measures pricing efficiency and isolates premium transcontinental margins."
        }
    ]
    
    kpi_width = Inches(3.64)
    kpi_gap = Inches(0.4)
    kpi_start_left = Inches(0.8)
    kpi_top = Inches(2.2)
    kpi_height = Inches(4.4)
    
    for idx, kpi in enumerate(kpi_cards):
        k_left = kpi_start_left + (idx * (kpi_width + kpi_gap))
        
        # Add Card Base Shape
        add_card(slide11, k_left, kpi_top, kpi_width, kpi_height, border_color=COLOR_PRIMARY)
        
        # Card header accent
        accent_bar = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, k_left, kpi_top, kpi_width, Inches(0.1))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLOR_ACCENT
        accent_bar.line.fill.background()
        
        # Content box
        k_box = slide11.shapes.add_textbox(k_left + Inches(0.15), kpi_top + Inches(0.25), kpi_width - Inches(0.3), kpi_height - Inches(0.4))
        k_tf = k_box.text_frame
        k_tf.word_wrap = True
        
        # KPI Number
        p_num = k_tf.paragraphs[0]
        p_num.text = kpi["num"]
        p_num.font.name = FONT_TITLE
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_ACCENT
        p_num.space_after = Pt(4)
        
        # KPI Title
        p_title = k_tf.add_paragraph()
        p_title.text = kpi["title"]
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_LIGHT
        p_title.space_after = Pt(14)
        
        # KPI Description
        p_desc = k_tf.add_paragraph()
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10.5)
        p_desc.line_spacing = 1.15
        
        # Highlight strong statements in description
        parts = kpi["desc"].split("**")
        for i, part in enumerate(parts):
            run = p_desc.add_run()
            run.text = part
            run.font.name = FONT_BODY
            run.font.size = Pt(10.5)
            run.font.color.rgb = COLOR_TEXT_LIGHT if i % 2 == 1 else COLOR_TEXT_MUTED
            if i % 2 == 1:
                run.font.bold = True
                run.font.color.rgb = COLOR_GOLD

    # =========================================================================
    # SLIDE 12: CONCLUSION & FUTURE STEPS (READY FOR TAKE OFF)
    # =========================================================================
    print("Building Slide 12: Conclusion & Next Steps...")
    slide12 = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide12)
    add_header(slide12, "11 / Next Steps", "Ready for Takeoff: Future Roadmap")
    
    # Left Column: Strategic Roadmap
    roadmap_bullets = [
        "**Airport Slot & Gate Audit**: Secure slot entries and gate allocations at highly congested, slot-controlled hubs (JFK, LGA, DCA).",
        "**Full-Year Demand Analysis**: Expand modeling to Q2-Q4 DOT databases to capture seasonal passenger pricing elasticity.",
        "**Competitor Share Modeling**: Assess capacity shares of dominant carriers (Delta, United, American) to optimize entry pricing.",
        "**Baggage Fee Elasticity**: Model incremental impact of bag pricing adjustments ($30 vs $35 vs $40) on load factor."
    ]
    add_bullet_list(slide12, Inches(0.8), Inches(1.8), Inches(6.5), Inches(5.0), roadmap_bullets)
    
    # Right Column: Big Final Graphic card
    add_card(slide12, Inches(7.8), Inches(1.8), Inches(4.733), Inches(5.0), border_color=COLOR_ACCENT)
    
    # Accent background image container
    logo_container = slide12.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(9.166), Inches(2.2), Inches(2.0), Inches(2.0))
    logo_container.fill.solid()
    logo_container.fill.fore_color.rgb = COLOR_PRIMARY
    logo_container.line.color.rgb = COLOR_ACCENT
    logo_container.line.width = Pt(2)
    
    # Inner diamond shape inside hexagon
    logo_star = slide12.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(9.666), Inches(2.7), Inches(1.0), Inches(1.0))
    logo_star.fill.solid()
    logo_star.fill.fore_color.rgb = COLOR_GOLD
    logo_star.line.fill.background()
    
    # Text block
    final_box = slide12.shapes.add_textbox(Inches(8.0), Inches(4.5), Inches(4.333), Inches(2.1))
    ftf = final_box.text_frame
    ftf.word_wrap = True
    
    p = ftf.paragraphs[0]
    p.text = "\"On time, for you.\""
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.space_after = Pt(10)
    
    p2 = ftf.add_paragraph()
    p2.text = "Operational excellence starts with data-driven network design. Our premium 5-route portfolio is financially primed to succeed."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_BODY
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.line_spacing = 1.2
    # Save the presentation
    filename = "Airline_Analysis_Presentation.pptx"
    print(f"Saving presentation as {filename}...")
    try:
        prs.save(filename)
    except PermissionError:
        fallback_filename = "Airline_Analysis_Presentation_New.pptx"
        print(f"WARNING: Permission denied when saving to {filename}. The file might be open in PowerPoint. Saving as {fallback_filename} instead...")
        prs.save(fallback_filename)
    print("PowerPoint presentation successfully created!")

if __name__ == "__main__":
    create_deck()
